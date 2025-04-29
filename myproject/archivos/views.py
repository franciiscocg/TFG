from rest_framework.views import APIView
from rest_framework.response import Response
from rest_framework import status
from rest_framework.permissions import IsAuthenticated
from .serializers import FileUploadSerializer, UploadedFileSerializer
import pdfplumber
from .models import UploadedFile
from django.shortcuts import get_object_or_404
from rest_framework.pagination import PageNumberPagination
from django.core.files.base import ContentFile
import os
from pptx import Presentation

class StandardResultsSetPagination(PageNumberPagination):
    page_size = 6  # 4 archivos por página
    page_size_query_param = 'page_size'
    max_page_size = 100

class FileUploadView(APIView):
    permission_classes = [IsAuthenticated]

    def post(self, request):
        serializer = FileUploadSerializer(data=request.data)
        if serializer.is_valid():
            serializer.save(user=request.user)
            return Response({"message": "Archivo subido con éxito"}, status=status.HTTP_201_CREATED)
        return Response(serializer.errors, status=status.HTTP_400_BAD_REQUEST)
 
class FileListView(APIView):
    permission_classes = [IsAuthenticated]
    pagination_class = StandardResultsSetPagination

    def get(self, request):
        files = UploadedFile.objects.filter(user=request.user)
        paginator = self.pagination_class()
        page = paginator.paginate_queryset(files, request)
        serializer = UploadedFileSerializer(page, many=True, context={'request': request})
        return paginator.get_paginated_response(serializer.data)

class FileDeleteView(APIView):
    permission_classes = [IsAuthenticated]

    def delete(self, request, file_id):
        file = get_object_or_404(UploadedFile, id=file_id, user=request.user)
        file.file.delete()
        if file.text_file:
            file.text_file.delete()
        file.delete()
        return Response({"message": "Archivo eliminado con éxito"}, status=status.HTTP_204_NO_CONTENT)

class ExtractTextView(APIView):
    permission_classes = [IsAuthenticated]

    def post(self, request, file_id):
        file_obj = get_object_or_404(UploadedFile, id=file_id, user=request.user)
        file_path = file_obj.file.path
        original_filename = file_obj.file.name

        # Obtener la extensión del archivo para determinar el tipo
        try:
            filename_base, file_extension = os.path.splitext(original_filename)
            file_extension = file_extension.lower() # Convertir a minúsculas para comparación
        except Exception as e:
             return Response(
                {"message": f"No se pudo determinar la extensión del archivo: {e}"},
                status=status.HTTP_400_BAD_REQUEST
            )
        extracted_text = ''
        error_message = None
        # --- Lógica de extracción basada en la extensión ---
        if file_extension == '.pdf':
            try:
                with pdfplumber.open(file_path) as pdf:
                    temp_text = []
                    for page in pdf.pages:
                        page_text = page.extract_text()
                        if page_text: # Añadir solo si hay texto en la página
                            temp_text.append(page_text)
                    extracted_text = '\n'.join(temp_text) # Unir páginas con salto de línea
            except Exception as e:
                error_message = f"Error al procesar el archivo PDF: {e}"

        elif file_extension == '.pptx':
            try:
                prs = Presentation(file_path)
                temp_text = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        # Verificar si la forma tiene un marco de texto y si contiene texto
                        if hasattr(shape, "text_frame") and shape.text_frame and shape.text_frame.text:
                            temp_text.append(shape.text_frame.text)
                extracted_text = '\n\n'.join(temp_text) # Unir texto de formas/diapositivas con doble salto
            except Exception as e:
                error_message = f"Error al procesar el archivo PowerPoint: {e}"

        else:
            # Si no es ni PDF ni PPTX, devolvemos un error
            return Response(
                {"message": f"Tipo de archivo no soportado: '{file_extension}'. Solo se admiten PDF y PPTX."},
                status=status.HTTP_400_BAD_REQUEST
            )

        # Si ocurrió un error durante la extracción específica del tipo de archivo
        if error_message:
             return Response({"message": error_message}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)

        # Verificar si se extrajo algo de texto útil
        if not extracted_text or not extracted_text.strip():
            return Response(
                {"message": f"No se pudo extraer contenido de texto del archivo ({file_extension}) o el archivo está vacío."},
                status=status.HTTP_400_BAD_REQUEST
            )

        output_filename_base = os.path.basename(filename_base)
        text_filename = f"{output_filename_base}.txt"

        try:
            file_obj.text_file.save(text_filename, ContentFile(extracted_text.encode('utf-8')), save=True)
            serializer = UploadedFileSerializer(file_obj, context={'request': request})
            return Response({
                "message": "Texto extraído y guardado con éxito",
                "file": serializer.data
            }, status=status.HTTP_200_OK)
        except Exception as e:
             return Response(
                 {"message": f"Error al guardar el archivo de texto extraído: {e}"},
                 status=status.HTTP_500_INTERNAL_SERVER_ERROR
             )


class ReadTextView(APIView):
    permission_classes = [IsAuthenticated]

    def get(self, request, file_id):
        file_obj = get_object_or_404(UploadedFile, id=file_id, user=request.user)
        if not file_obj.text_file:
            return Response({"message": "No hay texto extraído para este archivo"}, status=status.HTTP_400_BAD_REQUEST)
        try:
            with open(file_obj.text_file.path, 'r', encoding='utf-8') as f:
                text_content = f.read()
            return Response({
                "text": text_content,
                "extracted_data": file_obj.extracted_data if file_obj.extracted_data else {}
            }, status=status.HTTP_200_OK)
        except Exception as e:
            return Response({"message": f"Error al leer el archivo de texto: {str(e)}"}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)
        
class UserExtractedDataView(APIView):
    permission_classes = [IsAuthenticated]

    def get(self, request):
        # Obtener todos los archivos del usuario autenticado
        files = UploadedFile.objects.filter(user=request.user)
        
        # Filtrar solo los archivos que tienen extracted_data
        files_with_data = files.exclude(extracted_data__isnull=True).exclude(extracted_data={})
        
        if not files_with_data.exists():
            return Response({
                "message": "No se encontraron datos extraídos para este usuario",
                "count": 0,
                "data": []
            }, status=status.HTTP_200_OK)

        # Serializar los datos
        serializer = UploadedFileSerializer(files_with_data, many=True, context={'request': request})
        
        # Preparar la respuesta con solo los extracted_data
        extracted_data_list = [
            {
                "file_id": file["id"],
                "filename": file["file"].split('/')[-1],
                "extracted_data": file["extracted_data"]
            }
            for file in serializer.data if file["extracted_data"]
        ]

        return Response({
            "message": "Datos extraídos encontrados",
            "count": len(extracted_data_list),
            "data": extracted_data_list
        }, status=status.HTTP_200_OK)
        
class UpdateExtractedDataView(APIView):
    permission_classes = [IsAuthenticated]

    def put(self, request, file_id):
        file = get_object_or_404(UploadedFile, id=file_id, user=request.user)
        extracted_data = request.data.get('extracted_data', {})
        file.extracted_data = extracted_data
        file.save()
        return Response({'message': 'Datos actualizados con éxito'}, status=status.HTTP_200_OK)